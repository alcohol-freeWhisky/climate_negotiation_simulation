#!/usr/bin/env python3
"""Generate the LLM Agents in Scientific Research presentation (white + soft blue theme)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Palette  (white bg, low-saturation blues) ───────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG_WHITE    = RGBColor(0xFA, 0xFB, 0xFD)  # very faint blue-white
SOFT_BLUE   = RGBColor(0xE8, 0xEF, 0xF6)  # card backgrounds
MID_BLUE    = RGBColor(0xB0, 0xC9, 0xE0)  # borders, secondary
ACCENT      = RGBColor(0x3B, 0x7D, 0xD8)  # headings, accent bar
DARK_BLUE   = RGBColor(0x1E, 0x3A, 0x5F)  # primary text
NAVY        = RGBColor(0x0F, 0x1F, 0x3C)  # slide titles
GRAY        = RGBColor(0x6B, 0x7B, 0x8D)  # secondary text
LIGHT_GRAY  = RGBColor(0x9A, 0xA8, 0xB8)
GREEN_OK    = RGBColor(0x2E, 0x7D, 0x32)
RED_WARN    = RGBColor(0xC6, 0x28, 0x28)
HIGHLIGHT   = RGBColor(0x1A, 0x5C, 0xB0)  # darker accent for emphasis

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_bg(slide, color=BG_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def tb(slide, left, top, width, height, text, size=18, color=DARK_BLUE,
       bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def bullets(slide, left, top, width, height, items, size=17, color=DARK_BLUE,
            spacing=Pt(6), font="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    return txBox

def section_bar(slide, label):
    add_rect(slide, 0, 0, W, Inches(0.5), ACCENT)
    tb(slide, Inches(0.6), Inches(0.03), Inches(10), Inches(0.42),
       label, size=13, color=WHITE, bold=True)

def slide_title(slide, title, subtitle=None):
    tb(slide, Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.8),
       title, size=34, color=NAVY, bold=True)
    if subtitle:
        tb(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
           subtitle, size=17, color=GRAY)

def snum(slide, num, total=30):
    tb(slide, Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.35),
       f"{num} / {total}", size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

def card(slide, left, top, width, height, title, items, title_color=ACCENT,
         bg=SOFT_BLUE, border=MID_BLUE, text_size=16):
    add_rounded_rect(slide, left, top, width, height, bg, border)
    add_rect(slide, left + Inches(0.02), top + Inches(0.02), Inches(0.07), height - Inches(0.04), ACCENT)
    tb(slide, left + Inches(0.25), top + Inches(0.12), width - Inches(0.4), Inches(0.4),
       title, size=18, color=title_color, bold=True)
    bullets(slide, left + Inches(0.25), top + Inches(0.55), width - Inches(0.4),
            height - Inches(0.65), items, size=text_size, color=DARK_BLUE)

# Table helper
def make_table(slide, left, top, col_widths, headers, rows, header_bg=ACCENT,
               header_fg=WHITE, row_bg1=SOFT_BLUE, row_bg2=WHITE, text_size=14):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    total_w = sum(col_widths)
    row_h = Inches(0.55)

    # Header
    x = left
    for c, h in enumerate(headers):
        add_rect(slide, x, top, Inches(col_widths[c]), row_h, header_bg)
        tb(slide, x + Inches(0.08), top + Inches(0.08), Inches(col_widths[c]) - Inches(0.16),
           row_h, h, size=text_size, color=header_fg, bold=True, align=PP_ALIGN.CENTER)
        x += Inches(col_widths[c])

    # Data rows
    for r, row in enumerate(rows):
        y = top + row_h * (r + 1)
        bg = row_bg1 if r % 2 == 0 else row_bg2
        x = left
        for c, val in enumerate(row):
            add_rect(slide, x, y, Inches(col_widths[c]), row_h, bg, border_color=MID_BLUE)
            tb(slide, x + Inches(0.08), y + Inches(0.08), Inches(col_widths[c]) - Inches(0.16),
               row_h, val, size=text_size - 1, color=DARK_BLUE, align=PP_ALIGN.CENTER)
            x += Inches(col_widths[c])


TOTAL = 30

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
add_rect(sl, 0, Inches(2.2), W, Inches(3.2), SOFT_BLUE)
add_rect(sl, Inches(0.8), Inches(2.05), Inches(0.1), Inches(3.5), ACCENT)

tb(sl, Inches(1.3), Inches(2.4), Inches(11), Inches(1.0),
   "From Chatbots to Co-Pilots", size=44, color=NAVY, bold=True)
tb(sl, Inches(1.3), Inches(3.4), Inches(11), Inches(0.7),
   "LLM Agents & Autonomous Workflows in Scientific Research",
   size=24, color=ACCENT)
tb(sl, Inches(1.3), Inches(4.3), Inches(11), Inches(0.5),
   "Presenter  |  SERC Group Meeting  |  April 2026", size=16, color=GRAY)
snum(sl, 1, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Evolution
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 1  |  REDEFINING THE LLM AGENT")
slide_title(sl, "The Evolution of AI in Science")

boxes = [
    ("Generative AI", "Text & code\ngeneration"),
    ("Conversational\nLLMs", "ChatGPT-style\nQ&A interfaces"),
    ("Autonomous\nAgents", "Plan, execute,\nself-correct"),
    ("Multi-Agent\nSystems", "Virtual labs &\norchestrated teams"),
]
for i, (t, d) in enumerate(boxes):
    left = Inches(0.8 + i * 3.1)
    add_rounded_rect(sl, left, Inches(2.6), Inches(2.6), Inches(2.6), SOFT_BLUE, MID_BLUE)
    tb(sl, left + Inches(0.15), Inches(2.8), Inches(2.3), Inches(0.8),
       t, size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, left + Inches(0.15), Inches(3.7), Inches(2.3), Inches(1.0),
       d, size=15, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    if i < 3:
        tb(sl, Inches(3.45 + i * 3.1), Inches(3.4), Inches(0.6), Inches(0.6),
           "\u279C", size=30, color=ACCENT, align=PP_ALIGN.CENTER)

tb(sl, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.5),
   'Core thesis:  We are moving from AI as an "encyclopedia" to AI as an "execution engine."',
   size=18, color=HIGHLIGHT, bold=True)
snum(sl, 2, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What is an Agent
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 1  |  REDEFINING THE LLM AGENT")
slide_title(sl, 'What Exactly is an "Agent"?')

components = [
    ("Brain", "Foundation model (LLM)\nReasoning & generation"),
    ("Memory", "Vector DB (RAG) + conversation\ncontext + long-term recall"),
    ("Planning", "Task decomposition,\nChain-of-Thought, tree search"),
    ("Tools", "API calls, bash, Python\ninterpreter, file I/O"),
]
for i, (t, d) in enumerate(components):
    left = Inches(0.7 + i * 3.1)
    add_rounded_rect(sl, left, Inches(2.5), Inches(2.8), Inches(2.8), SOFT_BLUE, MID_BLUE)
    tb(sl, left + Inches(0.2), Inches(2.7), Inches(2.4), Inches(0.5),
       t, size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, left + Inches(0.2), Inches(3.4), Inches(2.4), Inches(1.4),
       d, size=16, color=DARK_BLUE, align=PP_ALIGN.CENTER)

tb(sl, Inches(0.8), Inches(5.7), Inches(11), Inches(0.5),
   "An agent = LLM + memory + planning + tool access, operating in a loop.",
   size=18, color=HIGHLIGHT, bold=True)
snum(sl, 3, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ReAct
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 1  |  REDEFINING THE LLM AGENT")
slide_title(sl, "The ReAct Paradigm (Reasoning + Acting)")

cycle = ["Observation", "Thought", "Action", "Observation\u2082"]
for i, label in enumerate(cycle):
    left = Inches(1.2 + i * 2.9)
    add_rounded_rect(sl, left, Inches(2.7), Inches(2.3), Inches(1.2), SOFT_BLUE, MID_BLUE)
    tb(sl, left, Inches(2.9), Inches(2.3), Inches(0.7),
       label, size=21, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    if i < 3:
        tb(sl, left + Inches(2.3), Inches(2.9), Inches(0.6), Inches(0.7),
           "\u279C", size=26, color=ACCENT, align=PP_ALIGN.CENTER)

bullets(sl, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5), [
    "The agent reasons about what it observes, takes an action, then re-evaluates.",
    "Why it matters for research: self-correction based on intermediate outputs.",
    "Example: Agent reads a solver log, spots infeasibility, adjusts constraints, re-runs.",
], size=18, color=DARK_BLUE)
snum(sl, 4, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MCP
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 1  |  REDEFINING THE LLM AGENT")
slide_title(sl, "MCP: The Universal Tool Protocol")

card(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.2),
     "Model Context Protocol (Anthropic)", [
         "A standard interface connecting LLMs to any external tool, data source, or service.",
         'Think of it as USB-C for AI agents \u2014 one protocol, infinite tools.',
         "Eliminates custom integration code per tool.",
         "Adopted by Claude Code, Cursor, Windsurf, and growing ecosystem.",
     ])
card(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.2),
     "Why MCP Matters for Labs", [
         "Connect agents to Slack, Google Drive, databases, HPC schedulers.",
         "Community MCP servers for SLURM, Jupyter, PostgreSQL already exist.",
         "Build once, reuse across any MCP-compatible client.",
         "Our lab can build a SLURM MCP server \u2192 every team member benefits.",
     ])
snum(sl, 5, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Evaluating Models
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 2  |  MODELS, PLATFORMS & PRICING")
slide_title(sl, "Evaluating Models: What Actually Matters for Research")

bullets(sl, Inches(0.8), Inches(2.3), Inches(11.5), Inches(4.5), [
    "Standard benchmarks (MMLU, HumanEval) do NOT reflect scientific utility.",
    "",
    "What actually matters for our workflows:",
    "\u2022  Tool-calling accuracy \u2014 Can the model reliably invoke APIs, bash, solvers?",
    "\u2022  Long-context retrieval \u2014 Can it find the right variable in a 5,000-line script?",
    "\u2022  Niche coding syntax \u2014 Julia, JuMP, Gurobi, GAMS, domain-specific DSLs.",
    "\u2022  Agentic capability \u2014 Can it plan, execute, and self-correct across multi-step tasks?",
    "\u2022  Cost efficiency \u2014 How much does a typical debugging session cost?",
], size=18, color=DARK_BLUE)
snum(sl, 6, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Claude Opus 4.7
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 2  |  MODELS, PLATFORMS & PRICING")
slide_title(sl, "Claude Opus 4.7 / Sonnet 4.6 (Anthropic)")

card(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5),
     "Strengths", [
         "#1 on SWE-bench \u2014 best agentic coding model.",
         "Lowest hallucination for math & optimization.",
         "1M token context, no long-context surcharge.",
         "Claude Code: terminal-native agent (reads/writes/commits).",
         "First-class MCP ecosystem.",
     ])
card(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5),
     "Pricing (API pay-per-token)", [
         "Opus 4.7:   $5 / M input, $25 / M output.",
         "Sonnet 4.6: $3 / M input, $15 / M output.",
         "Haiku 4.5:  $1 / M input, $5 / M output.",
         "Prompt caching: 90% savings on repeated context.",
         "Batch API: additional 50% off (async jobs).",
         "Combined savings up to 95%.",
     ])
snum(sl, 7, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GPT-5.4 / o4
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 2  |  MODELS, PLATFORMS & PRICING")
slide_title(sl, "GPT-5.4 / o3 / o4-mini (OpenAI)")

card(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5),
     "Strengths", [
         "GPT-5.4: flagship, 1.1M context, strong all-rounder.",
         "o3: advanced reasoning model for math-heavy tasks.",
         "o4-mini: most cost-effective reasoning model.",
         "Most stable function / API calling across the industry.",
         "Excellent multi-modal input (images, plots, diagrams).",
     ])
card(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5),
     "Pricing (API pay-per-token)", [
         "GPT-5.4:    $2.50 / M input, $15 / M output.",
         "GPT-4.1:    $2.00 / M input, $8 / M output.",
         "o3:         $2.00 / M input, $8 / M output.",
         "o4-mini:    $1.10 / M input, $4.40 / M output.",
         "GPT-5.4 doubles price above 272K context.",
         "Codex CLI: OpenAI's coding agent (like Claude Code).",
     ])
snum(sl, 8, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Gemini 3.1 Pro
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 2  |  MODELS, PLATFORMS & PRICING")
slide_title(sl, "Gemini 3.1 Pro / 2.5 Flash (Google)")

card(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5),
     "Strengths", [
         "Gemini 3.1 Pro: 77.1% ARC-AGI-2 (vs 31.1% for 3.0).",
         "45-80% gains on agentic benchmarks over predecessor.",
         "Gemini 2.5 Flash: 1M context at incredibly low cost.",
         "Best-in-class for high-volume, budget-sensitive workloads.",
         "Strong multimodal: images, video, audio natively.",
     ])
card(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5),
     "Pricing (API pay-per-token)", [
         "Gemini 3.1 Pro: $2.00 / M input, $12 / M output.",
         "  (doubles above 200K context)",
         "Gemini 2.5 Flash: $0.30 / M input, $2.50 / M output.",
         "  (cheapest frontier model available)",
         "Generous free tier for experimentation.",
         "Google One AI Premium: $20/mo consumer access.",
     ])
snum(sl, 9, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Open-Weight Models
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 2  |  MODELS, PLATFORMS & PRICING")
slide_title(sl, "Open-Weight Models: Free & Local Deployment")

card(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5),
     "Key Models", [
         "DeepSeek-V3 / R1 \u2014 top reasoning, fraction of API cost.",
         "Qwen 3 \u2014 strong multilingual + code, Apache 2.0 license.",
         "Llama 4 (Meta) \u2014 massive community, fine-tunable.",
         "Mistral Large \u2014 fast, efficient for tool use.",
     ])
card(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5),
     "Why Open-Weight for Labs?", [
         "Deploy on HPC (Great Lakes, ORCD) \u2014 data never leaves.",
         "No API costs at scale (GPU time only).",
         "Fine-tune on domain data (climate, energy, optimization).",
         "Ollama / vLLM for easy local serving.",
         "Full reproducibility: pin exact model weights.",
     ])
snum(sl, 10, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Model Comparison Table
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 2  |  MODELS, PLATFORMS & PRICING")
slide_title(sl, "Model Comparison at a Glance")

col_w = [2.0, 2.1, 2.1, 2.1, 2.1, 2.1]
headers = ["Dimension", "Claude 4.7", "GPT-5.4", "Gemini 3.1P", "Open-Weight", "Gemini Flash"]
rows = [
    ["Coding/Agents", "\u2605\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605", "\u2605\u2605\u2605"],
    ["Tool Calling", "\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605", "\u2605\u2605\u2605"],
    ["Reasoning", "\u2605\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605"],
    ["Context", "1M", "1.1M", "1M", "128K", "1M"],
    ["Vision", "\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605\u2605\u2605", "\u2605\u2605\u2605", "\u2605\u2605\u2605\u2605"],
    ["Local Deploy", "\u2717", "\u2717", "\u2717", "\u2713\u2713\u2713", "\u2717"],
    ["Input $/M", "$5.00", "$2.50", "$2.00", "Free", "$0.30"],
    ["Output $/M", "$25.00", "$15.00", "$12.00", "Free", "$2.50"],
]
make_table(sl, Inches(0.5), Inches(2.15), col_w, headers, rows)
snum(sl, 11, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Pay-by-Token Explainer
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 2  |  MODELS, PLATFORMS & PRICING")
slide_title(sl, "How Pay-by-Token Works")

card(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5),
     "The Token Economy", [
         "1 token \u2248 4 characters \u2248 0.75 words.",
         "You pay per token for both input AND output.",
         "Input: your prompt + system instructions + context.",
         "Output: the model's response.",
         "Longer conversations = higher cost (context accumulates).",
     ])
card(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5),
     "Real-World Cost Examples", [
         "A typical debugging session (5 rounds):",
         "  \u2248 50K tokens \u2192 ~$0.50-$2.00 depending on model.",
         "A full agentic workflow (30 min agent run):",
         "  \u2248 500K tokens \u2192 ~$5-$15.",
         "Prompt caching: reuse context \u2192 90% cheaper.",
         "Batch API: async jobs \u2192 50% cheaper.",
     ])
snum(sl, 12, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Integrated Platforms Overview
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 2  |  MODELS, PLATFORMS & PRICING")
slide_title(sl, "Integrated AI Coding Platforms")

card(sl, Inches(0.4), Inches(2.3), Inches(3.0), Inches(4.5),
     "Cursor", [
         "AI-native IDE (VS Code fork).",
         "Inline edits + chat + agent.",
         "Multi-model: Claude, GPT, Gemini.",
         "Students: 1yr Pro FREE (.edu).",
         "Pro $20/mo, Ultra $200/mo.",
     ], text_size=15)
card(sl, Inches(3.6), Inches(2.3), Inches(3.0), Inches(4.5),
     "GitHub Copilot", [
         "Deep VS Code / JetBrains integration.",
         "Agent mode (multi-file edits).",
         "Free tier: 2K completions/mo.",
         "Students: Pro FREE (edu pack).",
         "Pro $10/mo, Pro+ $39/mo.",
     ], text_size=15)
card(sl, Inches(6.8), Inches(2.3), Inches(3.0), Inches(4.5),
     "Windsurf", [
         "AI IDE with SWE-1 agent.",
         "Multi-model backend.",
         "Unlimited tab/inline edits.",
         "Students: ~$8-10/mo (.edu).",
         "Pro $20/mo, Max $40/mo.",
     ], text_size=15)
card(sl, Inches(10.0), Inches(2.3), Inches(3.0), Inches(4.5),
     "Claude Code", [
         "Terminal-native agent.",
         "Read/write/test/commit.",
         "MCP tool integration.",
         "Pro $20/mo or API key.",
         "Max $100/mo (5x capacity).",
     ], text_size=15)
snum(sl, 13, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Platform Comparison Table
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 2  |  MODELS, PLATFORMS & PRICING")
slide_title(sl, "Platform Pricing & Billing Model Comparison")

col_w = [2.0, 2.3, 2.3, 2.3, 2.3]
headers = ["Feature", "Cursor", "Copilot", "Windsurf", "Claude Code"]
rows = [
    ["Monthly Cost", "$20 (Pro)", "$10 (Pro)", "$20 (Pro)", "$20 (Pro)"],
    ["Student Price", "FREE 1yr", "FREE (edu)", "~$8/mo", "$20 (no disc.)"],
    ["Billing Model", "Sub + credits", "Sub + quota", "Sub + quota", "Sub or API key"],
    ["Multi-Model", "Yes (3+ LLMs)", "Yes (3+ LLMs)", "Yes (3+ LLMs)", "Claude only"],
    ["Agent Mode", "Yes", "Yes (new)", "Yes (SWE-1)", "Yes (native)"],
    ["Terminal Agent", "Limited", "CLI preview", "Limited", "Full (primary)"],
    ["MCP Support", "Yes", "No", "Yes", "Yes (native)"],
    ["Team/Lab Plan", "$40/seat/mo", "$19/seat/mo", "Contact sales", "$100/seat/mo"],
]
make_table(sl, Inches(0.6), Inches(2.1), col_w, headers, rows, text_size=13)
snum(sl, 14, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Lab Plan Recommendation
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 2  |  MODELS, PLATFORMS & PRICING")
slide_title(sl, "Recommendation: Which Plan Should Our Lab Subscribe?")

card(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.8),
     "Recommended Strategy", [
         "1. Everyone: claim Cursor Pro FREE (1yr, .edu email).",
         "2. Everyone: claim Copilot Pro FREE (GitHub edu pack).",
         "3. Lab plan: Anthropic API key (pay-per-token, shared).",
         "   \u2192 Use via Claude Code for heavy agentic work.",
         "   \u2192 Prompt caching + batch API = 90-95% savings.",
         "4. For HPC: deploy open-weight models (free, private).",
     ])
card(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.8),
     "Cost Estimate (10-person lab)", [
         "Cursor: $0/mo (student free tier).",
         "Copilot: $0/mo (student free tier).",
         "Claude API: ~$50-200/mo shared (usage-based).",
         "  (set budget alerts at $100, $200 thresholds)",
         "Open-weight on HPC: $0 (use existing GPU allocation).",
         "",
         "Total: ~$50-200/mo for the entire lab.",
     ])
snum(sl, 15, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — RAG for Literature
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 2  |  MODELS, PLATFORMS & PRICING")
slide_title(sl, "RAG for Literature Review & Paper Digestion")

card(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.2),
     "Ready-to-Use Tools", [
         "NotebookLM (Google) \u2014 upload PDFs, get Q&A + audio summaries.",
         "Elicit / Consensus \u2014 structured literature search + LLM synthesis.",
         "Semantic Scholar API + RAG pipeline \u2014 custom paper Q&A.",
         "Zotero + LLM plugins \u2014 chat with your entire library.",
     ])
card(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.2),
     "Custom RAG Pipeline", [
         "Chunk papers \u2192 embed \u2192 store in vector DB (Chroma).",
         "Query: LLM retrieves relevant chunks, then answers.",
         "Add citation tracking for reproducibility.",
         "Ideal for domain corpora (e.g., all IPCC reports).",
     ])
snum(sl, 16, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Why Single Agents Fail
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 3  |  MULTI-AGENT & MULTI-MODEL ORCHESTRATION")
slide_title(sl, "Why Single Agents Fail in Complex Research")

bullets(sl, Inches(0.8), Inches(2.4), Inches(11.5), Inches(4.5), [
    '\u2022  "Lost in the Middle" \u2014 context fills up; agent forgets primary objective.',
    "\u2022  Infinite loops \u2014 retries the same failing approach without escalating.",
    "\u2022  Role confusion \u2014 one agent can't be expert debugger AND careful analyst.",
    "\u2022  No separation of concerns \u2014 debugging pollutes the analysis context.",
    "\u2022  Single-model limitation \u2014 no one model is best at everything.",
    "",
    "Solution: multiple specialized agents, potentially using different models.",
], size=19, color=DARK_BLUE)
snum(sl, 17, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Supervisor-Worker
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 3  |  MULTI-AGENT & MULTI-MODEL ORCHESTRATION")
slide_title(sl, "The Supervisor-Worker Architecture")

# Supervisor
add_rounded_rect(sl, Inches(4.2), Inches(2.2), Inches(4.8), Inches(1.2), ACCENT, MID_BLUE)
tb(sl, Inches(4.2), Inches(2.35), Inches(4.8), Inches(0.9),
   "SUPERVISOR  (Opus 4.7)\nDecomposes task, assigns, reviews",
   size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

workers = [
    ("Data Engineer", "Sonnet 4.6", "Clean, transform,\nvalidate datasets"),
    ("Solver Agent", "GPT-5.4", "Formulate & run\noptimization"),
    ("SLURM Operator", "Haiku 4.5", "Write SBATCH,\nsubmit, monitor"),
    ("Analyst", "Gemini 3.1", "Interpret results,\ngenerate plots"),
]
for i, (title, model, desc) in enumerate(workers):
    left = Inches(0.6 + i * 3.15)
    add_rounded_rect(sl, left, Inches(4.5), Inches(2.8), Inches(2.3), SOFT_BLUE, MID_BLUE)
    tb(sl, left + Inches(0.1), Inches(4.6), Inches(2.6), Inches(0.4),
       title, size=17, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, left + Inches(0.1), Inches(5.0), Inches(2.6), Inches(0.35),
       f"({model})", size=13, color=GRAY, align=PP_ALIGN.CENTER)
    tb(sl, left + Inches(0.1), Inches(5.35), Inches(2.6), Inches(1.0),
       desc, size=14, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    # Arrow
    tb(sl, left + Inches(1.0), Inches(3.5), Inches(0.8), Inches(0.9),
       "\u2193", size=32, color=ACCENT, align=PP_ALIGN.CENTER)

snum(sl, 18, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Multi-Model Orchestration
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 3  |  MULTI-AGENT & MULTI-MODEL ORCHESTRATION")
slide_title(sl, "Multi-Model Orchestration: Best Model for Each Task")

card(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5),
     "Why Mix Models?", [
         "No single model dominates every dimension.",
         "Claude: best at coding & debugging agents.",
         "GPT-5.4: best at tool calling & vision.",
         "Gemini 3.1: best reasoning per dollar.",
         "Gemini Flash / Haiku: cheapest for simple routing tasks.",
         "Open-weight: free for HPC, private data.",
     ])
card(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5),
     "How to Implement", [
         "LangGraph: each node can use a different LLM.",
         "Router pattern: cheap model triages \u2192 expensive model executes.",
         "Cost optimization: use Haiku/Flash for 80% of tokens, Opus for 20%.",
         "Fallback chains: if Model A fails, try Model B.",
         "A/B testing: compare model outputs on same task.",
     ])
snum(sl, 19, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Frameworks
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 3  |  MULTI-AGENT & MULTI-MODEL ORCHESTRATION")
slide_title(sl, "Frameworks: LangGraph vs. CrewAI vs. Claude Agent SDK")

card(sl, Inches(0.4), Inches(2.3), Inches(3.9), Inches(4.2),
     "LangGraph", [
         "State-machine workflows.",
         "Explicit edges & conditions.",
         "Best for strict scientific pipelines.",
         "Multi-model: each node = different LLM.",
         'If solver_error \u2192 route to Debugger.',
     ])
card(sl, Inches(4.6), Inches(2.3), Inches(3.9), Inches(4.2),
     "CrewAI", [
         'Role-based "company" metaphor.',
         "Quick setup: define roles & goals.",
         "Best for exploratory workflows.",
         "Built-in delegation & memory.",
         "Easy to prototype in an afternoon.",
     ])
card(sl, Inches(8.8), Inches(2.3), Inches(3.9), Inches(4.2),
     "Claude Agent SDK", [
         "Anthropic's native multi-agent framework.",
         "First-class MCP tool support.",
         "Handoffs between specialized agents.",
         "Guardrails built in.",
         "Ideal for Claude-based pipelines.",
     ])
snum(sl, 20, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Code Snippet
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 3  |  MULTI-AGENT & MULTI-MODEL ORCHESTRATION")
slide_title(sl, "Example: LangGraph with Multi-Model Routing")

code = '''from langgraph.graph import StateGraph
from langchain_anthropic import ChatAnthropic
from langchain_openai import ChatOpenAI

# Different models for different tasks
solver_llm  = ChatAnthropic(model="claude-opus-4-7")    # best at code
analyst_llm = ChatOpenAI(model="gpt-5.4")               # best at vision
router_llm  = ChatAnthropic(model="claude-haiku-4-5")   # cheapest for routing

def check_convergence(state):
    if state["solver_status"] == "optimal":
        return "analyze"        # proceed to analysis (GPT-5.4)
    elif state["retry_count"] >= 3:
        return "human_review"   # escalate to human
    else:
        return "debug"          # route to debugger (Opus)

graph = StateGraph(SolverState)
graph.add_conditional_edges("solve", check_convergence)
workflow = graph.compile()'''

add_rounded_rect(sl, Inches(0.8), Inches(2.1), Inches(11.7), Inches(5.0), SOFT_BLUE, MID_BLUE)
tb(sl, Inches(1.1), Inches(2.3), Inches(11.1), Inches(4.6),
   code, size=14, color=DARK_BLUE, font="Courier New")
snum(sl, 21, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — HPC Bridge
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 4  |  HPC & LOCAL DEPLOYMENT")
slide_title(sl, "Bridging the LLM and the Cluster")

card(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5),
     "The Challenge", [
         "LLMs live in the cloud (API calls).",
         "Our data & heavy compute live on SLURM clusters.",
         "Sensitive data cannot leave the cluster.",
         "Latency: API round-trips add overhead.",
     ])
card(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5),
     "Solutions", [
         "Option A: SSH-tunneled agents (Claude Code / Aider over SSH).",
         "Option B: MCP servers running on login nodes.",
         "Option C: Local open-weight models via vLLM on GPU nodes.",
         "Option D: Hybrid \u2014 cloud API for reasoning, local for data I/O.",
     ])
snum(sl, 22, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Local Deployment How-To
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 4  |  HPC & LOCAL DEPLOYMENT")
slide_title(sl, "How to Deploy LLMs on HPC (Step by Step)")

bullets(sl, Inches(0.8), Inches(2.3), Inches(11.5), Inches(5.0), [
    "1.  Request a GPU node:  salloc --gres=gpu:a100:1 --mem=80G --time=4:00:00",
    "",
    "2.  Install vLLM (or Ollama) in a conda env:",
    "      conda create -n llm python=3.11 && pip install vllm",
    "",
    "3.  Download model weights (one-time):",
    "      huggingface-cli download deepseek-ai/DeepSeek-V3 --local-dir /scratch/models/",
    "",
    "4.  Launch the server (OpenAI-compatible API):",
    "      python -m vllm.entrypoints.openai.api_server --model /scratch/models/DeepSeek-V3",
    "",
    "5.  Connect your agent (any OpenAI SDK client works):",
    "      export OPENAI_BASE_URL=http://gpu-node:8000/v1",
    "",
    "6.  SSH tunnel from laptop:  ssh -L 8000:gpu-node:8000 login.hpc.edu",
], size=16, color=DARK_BLUE, font="Calibri")
snum(sl, 23, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — SLURM Agent
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 4  |  HPC & LOCAL DEPLOYMENT")
slide_title(sl, "Automated Job Scheduling: The SLURM Agent")

steps = [
    ("1. Read", "Parse experiment\nscripts (Py/Julia)"),
    ("2. Estimate", "Predict resource\nneeds (GPU, RAM)"),
    ("3. Write", "Generate #SBATCH\ndirectives"),
    ("4. Submit", "sbatch + monitor\nsqueue output"),
    ("5. React", "Parse logs, fix\nerrors, resubmit"),
]
for i, (title, desc) in enumerate(steps):
    left = Inches(0.4 + i * 2.55)
    add_rounded_rect(sl, left, Inches(2.6), Inches(2.3), Inches(2.2), SOFT_BLUE, MID_BLUE)
    tb(sl, left + Inches(0.1), Inches(2.7), Inches(2.1), Inches(0.5),
       title, size=19, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, left + Inches(0.1), Inches(3.3), Inches(2.1), Inches(1.2),
       desc, size=15, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    if i < 4:
        tb(sl, left + Inches(2.25), Inches(3.2), Inches(0.4), Inches(0.6),
           "\u279C", size=22, color=ACCENT, align=PP_ALIGN.CENTER)

tb(sl, Inches(0.8), Inches(5.3), Inches(11), Inches(0.5),
   "End-to-end: from experiment code to monitored cluster job, zero manual SBATCH writing.",
   size=17, color=HIGHLIGHT, bold=True)
snum(sl, 24, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Auto-Healing Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 4  |  HPC & LOCAL DEPLOYMENT")
slide_title(sl, 'The "Auto-Healing" Pipeline')

bullets(sl, Inches(0.8), Inches(2.3), Inches(11.5), Inches(4.5), [
    "Scenario: A parallel array job crashes with CUDA Out-of-Memory (OOM).",
    "",
    '1.  Monitor Agent detects non-zero exit code, reads the .err log.',
    '2.  Diagnosis: "RuntimeError: CUDA out of memory. Tried to allocate 2.4 GiB."',
    "3.  Agent edits batch script: increases --mem or --gres=gpu allocation.",
    "4.  Agent resubmits job, logs the change for the researcher.",
    "5.  On success, proceeds to the analysis stage automatically.",
    "",
    "Key: Every auto-fix is logged. Production runs require human sign-off.",
], size=18, color=DARK_BLUE)
snum(sl, 25, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Energy Modeling Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 5  |  DIRECT RESEARCH APPLICATIONS")
slide_title(sl, "Automated Energy Modeling & Downscaling Pipeline")

card(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.2),
     "Data Prep Agent", [
         "Automates ERA5 / CMIP6 extraction.",
         "Targets exact params: rsds, huss, tas, pr.",
         "Handles CDS API auth, spatial subsetting, temporal alignment.",
         "Output: analysis-ready NetCDF / Parquet.",
     ])
card(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.2),
     "Solver Agent", [
         "Translates math formulation \u2192 solver code.",
         "Supports HiGHS / Gurobi / CPLEX.",
         "Validates constraint dimensions automatically.",
         "Runs sensitivity analysis on key parameters.",
     ])
snum(sl, 26, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — Climate Negotiation Sim
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 5  |  DIRECT RESEARCH APPLICATIONS")
slide_title(sl, "Agents as Research Subjects: Climate Negotiation Sim")

card(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.2),
     "Multi-Agent Game Theory", [
         "LLMs act as heterogeneous national actors.",
         "Each agent: unique priorities, constraints, GDP, emissions.",
         "Agents negotiate, form coalitions, propose treaties.",
         "Beyond code gen \u2014 LLMs as subjects of social science.",
     ])
card(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.2),
     "Research Questions", [
         "Can LLM agents reach Pareto-efficient agreements?",
         "How do injected biases shift negotiation outcomes?",
         "Does agent architecture affect cooperation rates?",
         "Can we model COP dynamics computationally?",
     ])
snum(sl, 27, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Pitfalls
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 6  |  PITFALLS, REALITIES & GETTING STARTED")
slide_title(sl, "The Dangers of Autonomous Workflows")

card(sl, Inches(0.4), Inches(2.3), Inches(3.9), Inches(4.2),
     "Hallucinated Code", [
         'Solver returns "optimal" for the WRONG formulation.',
         "Subtle sign errors in constraints.",
         "Confident but incorrect API usage.",
     ])
card(sl, Inches(4.6), Inches(2.3), Inches(3.9), Inches(4.2),
     "Reproducibility", [
         "Dynamic bug fixes = untraceable changes.",
         "Non-deterministic outputs across runs.",
         '"It worked yesterday" syndrome.',
     ])
card(sl, Inches(8.8), Inches(2.3), Inches(3.9), Inches(4.2),
     "Cost & Token Burn", [
         "Agent loops can burn $50+ in API calls.",
         "Long context = expensive windows.",
         "Budget guardrails are essential.",
     ])
snum(sl, 28, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 29 — Human in the Loop
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 6  |  PITFALLS, REALITIES & GETTING STARTED")
slide_title(sl, 'The "Human-in-the-Loop" Mandate')

card(sl, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5),
     "Design Principles", [
         "Agents are co-pilots, NOT principal investigators.",
         "Hard-code approval checkpoints in LangGraph.",
         "No destructive commands without human sign-off.",
         "All agent actions logged to git / W&B / LangSmith.",
     ])
card(sl, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5),
     "Reproducibility Solutions", [
         "Git-based agent traces (auto-commit each step).",
         "temperature=0 + fixed seeds for determinism.",
         "LangSmith / W&B for full agent execution traces.",
         "Version-pin model IDs (not just 'gpt-5').",
     ])
snum(sl, 29, TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 30 — Next Steps + Q&A
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
section_bar(sl, "SECTION 6  |  GETTING STARTED")
slide_title(sl, "Actionable Next Steps for Our Lab")

bullets(sl, Inches(0.8), Inches(2.3), Inches(11.5), Inches(3.5), [
    "1.  Everyone: claim Cursor Pro (free, .edu) and Copilot Pro (free, GitHub edu pack).",
    "2.  Set up a shared Anthropic API key with budget alerts ($100/$200 thresholds).",
    "3.  Deploy one open-weight model on Great Lakes (vLLM + DeepSeek-V3).",
    "4.  Build a shared RAG pipeline over our group's paper corpus (NotebookLM to start).",
    "5.  Prototype one multi-agent workflow in LangGraph for a real experiment.",
    "6.  Establish lab conventions: logging, checkpoints, cost budgets, model version pins.",
], size=18, color=DARK_BLUE)

add_rounded_rect(sl, Inches(3.8), Inches(5.8), Inches(5.8), Inches(1.0), ACCENT)
tb(sl, Inches(3.8), Inches(5.9), Inches(5.8), Inches(0.8),
   "Q & A", size=34, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
snum(sl, 30, TOTAL)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "/Users/ziqiwei/Desktop/2026spring/SERC/climate-negotiation-sim/LLM_Agents_Scientific_Research.pptx"
prs.save(out)
print(f"Saved {TOTAL} slides to {out}")
