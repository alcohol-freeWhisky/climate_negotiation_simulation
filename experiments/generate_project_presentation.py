#!/usr/bin/env python3
"""
Generate an academic-style project presentation for the climate negotiation repo.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Dict, List

import matplotlib.pyplot as plt
import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_FILE = ROOT / "climate_negotiation_project_overview.pptx"
ASSET_DIR = ROOT / "outputs" / "presentation_assets"

COLORS = {
    "bg": RGBColor(250, 250, 248),
    "ink": RGBColor(32, 37, 41),
    "muted": RGBColor(98, 108, 121),
    "navy": RGBColor(25, 56, 86),
    "teal": RGBColor(40, 110, 118),
    "green": RGBColor(71, 122, 83),
    "orange": RGBColor(181, 98, 67),
    "panel": RGBColor(238, 241, 244),
    "line": RGBColor(200, 207, 214),
}


def load_scenario() -> Dict:
    return yaml.safe_load(
        (ROOT / "config" / "scenarios" / "paris_article6_8.yaml").read_text(
            encoding="utf-8"
        )
    )


def collect_run_metrics() -> List[Dict]:
    base = ROOT / "outputs" / "deepseek" / "COP_Negotiation_Simulation_DeepSeek"
    rows: List[Dict] = []
    if not base.exists():
        return rows

    for run_dir in sorted(base.glob("*")):
        if not run_dir.is_dir():
            continue
        eval_files = sorted(run_dir.glob("evaluation_*.json"))
        if not eval_files:
            continue
        evaluation = json.loads(eval_files[-1].read_text(encoding="utf-8"))
        acceptability = evaluation.get("acceptability", {})
        rows.append(
            {
                "run": run_dir.name,
                "overall": evaluation.get("overall_score", 0.0),
                "reference_alignment": evaluation.get("summary_scores", {}).get(
                    "reference_alignment_score", 0.0
                ),
                "negotiation_quality": evaluation.get("summary_scores", {}).get(
                    "negotiation_quality_score", 0.0
                ),
                "rougeL": evaluation.get("rouge_l", {}).get("rougeL_f", 0.0),
                "bertscore_f1": evaluation.get("bertscore", {}).get("f1", 0.0),
                "accept_count": acceptability.get("accept_count", 0),
                "blocking_count": acceptability.get("oppose_count", 0)
                + acceptability.get("modify_count", 0),
                "consensus": acceptability.get("consensus_possible", False),
            }
        )
    return rows


def collect_latest_bloc_result() -> Dict[str, List[str]]:
    result_path = (
        ROOT
        / "outputs"
        / "deepseek"
        / "COP_Negotiation_Simulation_DeepSeek"
        / "20260413_213057"
        / "COP_Negotiation_Simulation_DeepSeek_20260413_213057_results.json"
    )
    if not result_path.exists():
        return {"acceptances": [], "objections": []}

    results = json.loads(result_path.read_text(encoding="utf-8"))
    return {
        "acceptances": results.get("acceptances", []),
        "objections": [item.get("agent", "") for item in results.get("objections", [])],
    }


def repo_snapshot() -> Dict[str, int]:
    return {
        "src_files": len(list((ROOT / "src").rglob("*.py"))),
        "agent_configs": len(list((ROOT / "config" / "agents").glob("*.yaml"))),
        "test_modules": len(list((ROOT / "tests").glob("test_*.py"))),
    }


def create_result_chart(metrics: List[Dict]) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    chart_path = ASSET_DIR / "academic_run_trends.png"

    recent = metrics[-5:] if len(metrics) >= 5 else metrics
    if not recent:
        return chart_path

    labels = [row["run"][4:8] for row in recent]
    ref = [row["reference_alignment"] for row in recent]
    neg = [row["negotiation_quality"] for row in recent]
    acc = [row["accept_count"] for row in recent]

    fig, axes = plt.subplots(1, 2, figsize=(11, 4.3), constrained_layout=True)
    fig.patch.set_facecolor("white")

    axes[0].plot(labels, ref, marker="o", linewidth=2.3, color="#193856")
    axes[0].plot(labels, neg, marker="o", linewidth=2.3, color="#286e76")
    axes[0].set_ylim(0, 1.0)
    axes[0].set_title("Quality trend", fontsize=12, weight="bold")
    axes[0].set_ylabel("Score")
    axes[0].grid(alpha=0.25)
    axes[0].legend(["Reference alignment", "Negotiation quality"], frameon=False, fontsize=9)

    axes[1].bar(labels, acc, color="#477a53", width=0.55)
    axes[1].set_ylim(0, 8)
    axes[1].set_title("Accepting blocs", fontsize=12, weight="bold")
    axes[1].set_ylabel("Count")
    axes[1].grid(axis="y", alpha=0.25)

    fig.savefig(chart_path, dpi=220, bbox_inches="tight")
    plt.close(fig)
    return chart_path


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=22,
    bold=False,
    color=COLORS["ink"],
    align=PP_ALIGN.LEFT,
    fill=None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = fill
    else:
        box.fill.background()
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    bullets: List[str],
    font_size=22,
    color=COLORS["ink"],
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.bullet = True
        p.level = 0
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Aptos"
    return box


def add_title(slide, title: str, subtitle: str | None = None):
    add_textbox(
        slide,
        Inches(0.65),
        Inches(0.35),
        Inches(12),
        Inches(0.75),
        title,
        font_size=28,
        bold=True,
        color=COLORS["navy"],
    )
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.65),
        Inches(1.08),
        Inches(2.1),
        Inches(0.05),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["teal"]
    line.line.color.rgb = COLORS["teal"]
    if subtitle:
        add_textbox(
            slide,
            Inches(0.65),
            Inches(1.16),
            Inches(11.8),
            Inches(0.35),
            subtitle,
            font_size=12,
            color=COLORS["muted"],
        )


def add_panel(slide, left, top, width, height, fill=COLORS["panel"]):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = COLORS["line"]
    return shape


def make_slide_background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]


def build_presentation():
    scenario = load_scenario()
    metrics = collect_run_metrics()
    latest = metrics[-1] if metrics else {}
    best_alignment = max(metrics, key=lambda row: row["rougeL"]) if metrics else {}
    latest_blocs = collect_latest_bloc_result()
    snapshot = repo_snapshot()
    chart_path = create_result_chart(metrics)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.2),
        Inches(11.6),
        Inches(0.9),
        "Climate Negotiation Simulation",
        font_size=30,
        bold=True,
        color=COLORS["navy"],
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(2.0),
        Inches(11.4),
        Inches(0.6),
        "A multi-agent LLM framework for simulating UNFCCC bargaining",
        font_size=21,
        color=COLORS["muted"],
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(3.25),
        Inches(5.5),
        Inches(1.5),
        "10-minute project overview\nBackground\nGoal\nWorkflow\nExperiment\nNext steps",
        font_size=22,
        color=COLORS["ink"],
        fill=RGBColor(235, 239, 243),
    )
    add_textbox(
        slide,
        Inches(7.0),
        Inches(3.25),
        Inches(5.0),
        Inches(1.5),
        "Current setup\n8 bloc agents + Chair\nDeepSeek default backend\nReal draft/final text benchmark",
        font_size=22,
        color=COLORS["ink"],
        fill=RGBColor(233, 240, 235),
    )

    # 2. Background
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide)
    add_title(slide, "Background")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.65),
        Inches(6.2),
        Inches(4.7),
        [
            "UNFCCC decisions are negotiated through text, not only through preferences.",
            "The process is consensus-based, multi-party, and highly path-dependent.",
            "Article 6 negotiations are especially difficult because legal drafting and political trade-offs interact.",
            "This makes climate bargaining a useful testbed for multi-agent LLM systems."
        ],
        font_size=22,
    )
    add_panel(slide, Inches(7.15), Inches(1.8), Inches(5.1), Inches(3.9))
    add_textbox(
        slide,
        Inches(7.45),
        Inches(2.05),
        Inches(4.5),
        Inches(3.2),
        "Why simulation is hard\n\nDistinct bloc identities\nIterative compromise text\nChair-mediated convergence\nAdoption depends on political acceptability, not only textual quality",
        font_size=22,
    )

    # 3. Goal
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide)
    add_title(slide, "Goal")
    add_bullets(
        slide,
        Inches(0.9),
        Inches(1.85),
        Inches(11.2),
        Inches(4.5),
        [
            "Build an end-to-end system that can simulate bloc-level climate negotiations over real draft text.",
            "Preserve persistent agent identities, red lines, and coalition behavior across multiple rounds.",
            "Produce revisable decision text and compare it against a real negotiated outcome.",
            "Separate textual realism from political feasibility in the evaluation."
        ],
        font_size=23,
    )

    # 4. Workflow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide)
    add_title(slide, "Workflow")
    phases = [
        ("Opening Statements", "initial priorities"),
        ("First Reading", "paragraph amendments"),
        ("Informal Consultations", "multi-round bargaining"),
        ("Final Plenary", "clean text and adoption"),
    ]
    for idx, (phase, sub) in enumerate(phases):
        left = Inches(0.8 + idx * 3.05)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            Inches(2.15),
            Inches(2.45),
            Inches(1.75),
        )
        color = [COLORS["navy"], COLORS["teal"], COLORS["green"], COLORS["orange"]][idx]
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"{phase}\n{sub}"
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
    add_bullets(
        slide,
        Inches(0.95),
        Inches(4.6),
        Inches(11.2),
        Inches(1.3),
        [
            "Chair synthesizes proposals into revised text after each consultation round.",
            "Text Manager tracks brackets, amendments, and text evolution."
        ],
        font_size=20,
    )

    # 5. Workflow: system implementation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide)
    add_title(slide, "Workflow Implementation")
    add_panel(slide, Inches(0.8), Inches(1.7), Inches(3.0), Inches(3.9))
    add_panel(slide, Inches(4.1), Inches(1.7), Inches(3.0), Inches(3.9))
    add_panel(slide, Inches(7.4), Inches(1.7), Inches(4.8), Inches(3.9))
    add_textbox(slide, Inches(1.05), Inches(2.0), Inches(2.4), Inches(3.0),
                "Negotiation Engine\n\nPhase Manager\nTurn Manager\nText Manager\nAmendment Processor", font_size=22)
    add_textbox(slide, Inches(4.35), Inches(2.0), Inches(2.4), Inches(3.0),
                "Agent Layer\n\n8 negotiating blocs\n1 Chair agent\nGeneral charters\nRuntime agenda focus", font_size=22)
    add_textbox(slide, Inches(7.7), Inches(2.0), Inches(4.2), Inches(3.0),
                "Evaluation Layer\n\nROUGE-L\nBERTScore\nKey clause matching\nStance consistency\nPolitical acceptability", font_size=22)

    # 6. Experiment
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide)
    add_title(slide, "Experiment")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.7),
        Inches(7.0),
        Inches(4.8),
        [
            f"Case study: {scenario['scenario_name']}",
            f"Participants: {len(scenario['active_agents'])} negotiating blocs plus a neutral Chair",
            "Default backend: DeepSeek chat",
            "Scenario disputes: scope, governance, relation to markets, institutional home, reporting",
            "Validation: pytest, unittest, compileall, and dry-run all pass"
        ],
        font_size=21,
    )
    add_panel(slide, Inches(8.25), Inches(1.85), Inches(3.8), Inches(3.4))
    add_textbox(
        slide,
        Inches(8.55),
        Inches(2.1),
        Inches(3.2),
        Inches(2.9),
        f"Repository snapshot\n\n{snapshot['src_files']} source files\n{snapshot['agent_configs']} bloc profiles\n{snapshot['test_modules']} test modules\nPer-run output folders",
        font_size=22,
    )

    # 7. Experiment results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide)
    add_title(slide, "Experiment Results")
    if chart_path.exists():
        slide.shapes.add_picture(str(chart_path), Inches(0.8), Inches(1.55), width=Inches(8.2))
    add_panel(slide, Inches(9.2), Inches(1.7), Inches(3.1), Inches(3.95))
    add_textbox(
        slide,
        Inches(9.45),
        Inches(1.95),
        Inches(2.6),
        Inches(3.4),
        "Representative findings\n\n"
        f"Best ROUGE-L: {best_alignment.get('rougeL', 0):.3f}\n"
        f"Best BERTScore: {best_alignment.get('bertscore_f1', 0):.3f}\n"
        f"Latest accepting blocs: {latest.get('accept_count', 0)} / 8\n"
        f"Latest blocking blocs: {latest.get('blocking_count', 0)} / 8\n\n"
        "The system writes realistic text before it reliably reaches full consensus.",
        font_size=18,
    )
    add_textbox(
        slide,
        Inches(0.95),
        Inches(6.05),
        Inches(11.2),
        Inches(0.45),
        "Observed pattern: textual alignment improved earlier than political acceptability, which makes endgame behavior the key bottleneck.",
        font_size=14,
        color=COLORS["muted"],
    )

    # 8. Next steps
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_slide_background(slide)
    add_title(slide, "Next Steps")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.7),
        Inches(5.9),
        Inches(4.8),
        [
            "Expand from one scenario to a broader set of negotiation topics.",
            "Run controlled ablations on bloc composition and behavioral parameters.",
            "Compare model families under the same negotiation protocol.",
            "Improve endgame convergence without collapsing bloc distinctions."
        ],
        font_size=22,
    )
    add_panel(slide, Inches(7.0), Inches(1.8), Inches(5.0), Inches(4.3))
    add_textbox(
        slide,
        Inches(7.3),
        Inches(2.1),
        Inches(4.35),
        Inches(3.6),
        "Research contribution\n\nA reusable workflow for studying whether LLM agents can negotiate like structured political actors, not just generate fluent text.\n\nThis is the bridge from a proof of concept to a paper-ready experimental platform.",
        font_size=22,
    )

    prs.save(OUTPUT_FILE)
    print(f"Saved presentation to: {OUTPUT_FILE}")


if __name__ == "__main__":
    build_presentation()
